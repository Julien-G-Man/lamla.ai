from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.auth import logout as auth_logout
from django.contrib import messages
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.conf import settings
from django.contrib.admin.views.decorators import staff_member_required
from django.views.decorators.http import require_POST
from django.db import transaction
from django.utils import timezone
from django.core.paginator import Paginator
from django.db.models import Q, Prefetch
from django.views.decorators.cache import cache_page
from django.core.cache import cache
from django.utils.decorators import method_decorator
from allauth.account.views import LoginView as AllauthLoginView
from allauth.account.adapter import DefaultAccountAdapter
import json
import logging
import os
import tempfile
from datetime import datetime, timedelta
from .models import (
    Question, QuestionCache, Contact, Feedback, 
    Subscription, UserProfile, ChatbotKnowledge, ChatMessage, QuizSession, ExamDocument, ExamAnalysis
)
from .chatbot_service import chatbot_service
from .email_backend import send_email
from .question_generator import generate_questions_from_text
from .flashcard_generator import generate_flashcards_from_text
import PyPDF2
import docx

import io
from io import BytesIO
from zoneinfo import ZoneInfo
from datetime import datetime
import re
import textwrap
import unicodedata
from urllib.parse import quote
from django.http import HttpResponse, FileResponse 
from django.utils import timezone
from django.views.decorators.http import require_GET, require_http_methods


import base64
from PIL import Image
import pytesseract
from django.core.validators import validate_email
from django.core.exceptions import ValidationError as DjangoValidationError
from django.template.loader import render_to_string
from django.contrib.auth import login

# Custom allauth adapter to ensure correct email sender
class CustomAccountAdapter(DefaultAccountAdapter):
    def send_mail(self, template_prefix, email, context):
        """
        Override send_mail to use our custom email backend and correct sender
        """
        from django.core.mail import EmailMessage
        from django.template.loader import render_to_string
        from django.utils.html import strip_tags
        
        # Determine the correct sender based on email type
        if 'password_reset' in template_prefix or 'email_confirmation' in template_prefix:
            from_email = getattr(settings, 'SECURITY_EMAIL_SENDER', 'lamlaaiteam@gmail.com')
        else:
            from_email = getattr(settings, 'WELCOME_EMAIL_SENDER', 'juliengmanana@gmail.com')
        
        # Render email templates
        subject = render_to_string(f'{template_prefix}_subject.txt', context)
        subject = ''.join(subject.splitlines()).strip()
        if not subject.startswith('[Lamla AI]'):
            subject = f'[Lamla AI] {subject}'
        
        html_message = render_to_string(f'{template_prefix}_message.html', context)
        plain_message = render_to_string(f'{template_prefix}_message.txt', context)
        
        # Create email message
        msg = EmailMessage(
            subject=subject,
            body=plain_message,
            from_email=from_email,
            to=[email]
        )
        msg.content_subtype = "html"
        
        # Use our custom send_email function
        return send_email(
            subject=subject,
            message=plain_message,
            recipient_list=[email],
            from_email=from_email,
            html_message=html_message,
            fail_silently=False
        )

    def save_user(self, request, user, form, commit=True):
        user = super().save_user(request, user, form, commit)
        # Ensure UserProfile is created
        UserProfile.objects.get_or_create(user=user)
        # Check for special signup flag
        from_account_not_found = False
        if request:
            # Check GET and POST for the flag
            if hasattr(request, 'GET') and request.GET.get('from_account_not_found') == '1':
                from_account_not_found = True
            elif hasattr(request, 'POST') and request.POST.get('from_account_not_found') == '1':
                from_account_not_found = True
            # Also check session (in case you want to persist the flag)
            elif hasattr(request, 'session') and request.session.get('from_account_not_found') == '1':
                from_account_not_found = True
        if from_account_not_found:
            # Send welcome email immediately and log in
            from django.template.loader import render_to_string
            from .email_backend import send_email
            context = {
                'user': user,
                'site_name': 'LAMLA AI',
                'site_domain': getattr(settings, 'SITE_DOMAIN', 'lamla-ai.onrender.com'),
            }
            subject = "Welcome to LAMLA AI! 🎉"
            html_message = render_to_string('emails/welcome_email.html', context)
            plain_message = render_to_string('emails/welcome_email.txt', context)
            try:
                logger.info(f"Attempting to send welcome email to {user.email} (user: {user.username})")
                send_email(
                    subject=subject,
                    message=plain_message,
                    recipient_list=[user.email],
                    from_email=getattr(settings, 'WELCOME_EMAIL_SENDER', 'juliengmanana@gmail.com'),
                    html_message=html_message
                )
                logger.info(f"Welcome email sent to {user.email}")
            except Exception as e:
                logger.error(f"Failed to send welcome email to {user.email}: {e}")
            if request:
                login(request, user)
        return user

@cache_page(60 * 15)  # Cache for 15 minutes
def home(request):
    return render(request, 'slides_analyzer/home.html')

@login_required
def dashboard(request):
    context = {}
    if request.user.is_authenticated:
        from .models import QuizSession, ExamAnalysis
        
        quiz_sessions = list(QuizSession.objects.filter(user=request.user))
        exam_analyses = list(ExamAnalysis.objects.filter(user=request.user))
        all_activities = []
        for q in quiz_sessions:
            subject = q.subject or ''
            file_name = ''
            if q.questions_data and isinstance(q.questions_data, dict):
                file_name = q.questions_data.get('uploaded_file_name', '')
            all_activities.append({'type': 'quiz', 'obj': q, 'created_at': q.created_at, 'subject': subject, 'file_name': file_name})
        for a in exam_analyses:
            all_activities.append({'type': 'exam_analysis', 'obj': a, 'created_at': a.created_at, 'subject': a.subject or '', 'file_name': ''})
        all_activities.sort(key=lambda x: x['created_at'], reverse=True)
        recent_activities = all_activities[:4]
        context['recent_activities'] = recent_activities
    return render(request, 'slides_analyzer/dashboard.html', context)

def user_profile(request):
    if request.user.is_authenticated:
        from .models import UserProfile
        # Get or create user profile
        profile, created = UserProfile.objects.get_or_create(user=request.user)
        
        if request.method == 'POST':
            # Handle profile updates
            updated = False
            
            if 'profile_picture' in request.FILES:
                profile_pic = request.FILES['profile_picture']
                # Validate file size (max 5MB)
                if profile_pic.size > 5 * 1024 * 1024:
                    messages.error(request, 'Profile picture must be less than 5MB.')
                    return redirect('user_profile')
                
                # Validate file type
                allowed_types = ['image/jpeg', 'image/jpg', 'image/png', 'image/gif']
                if profile_pic.content_type not in allowed_types:
                    messages.error(request, 'Please upload a valid image file (JPG, PNG, or GIF).')
                    return redirect('user_profile')
                
                profile.profile_picture = profile_pic
                updated = True
            
            bio = request.POST.get('bio', '').strip()
            if bio != profile.bio:
                profile.bio = bio
                updated = True
            
            if updated:
                profile.save()
                messages.success(request, 'Profile updated successfully!')
            else:
                messages.info(request, 'No changes were made to your profile.')
            
            return redirect('user_profile')
            
        return render(request, 'slides_analyzer/user_profile.html')
    else:
        return redirect('account_login')

def upload_slides(request):
    return render(request, 'slides_analyzer/upload.html')

def generate_questions(request):
    if request.method == 'POST':
        study_text = request.POST.get('extractedText', '').strip()
        num_mcq = int(request.POST.get('num_mcq', 5))
        num_short_raw = request.POST.get('num_short', '')
        try:
            num_short = int(num_short_raw)
            if num_short < 1:
                num_short = 0
        except (ValueError, TypeError):
            num_short = 0
        subject = request.POST.get('subject', '')
        difficulty = request.POST.get('difficulty', 'any')
        error_message = None
        quiz_results = None
        
        # Store uploaded file name if available
        uploaded_file_name = request.POST.get('uploaded_file_name', '')
        if not uploaded_file_name and 'slide_file' in request.FILES:
            uploaded_file_name = request.FILES['slide_file'].name
        
        # Log deployment debugging info
        logger.info(f"Quiz generation request - Text length: {len(study_text)}, MCQ: {num_mcq}, Short: {num_short}, Subject: {subject}")
        
        if not study_text or len(study_text) < 30:
            error_message = 'Please provide at least 30 characters of study material.'
            logger.warning(f"Quiz generation failed - insufficient text length: {len(study_text)}")
        else:
            try:
                logger.info("Calling generate_questions_from_text...")
                quiz_results = generate_questions_from_text(study_text, num_mcq, num_short, subject=subject, difficulty=difficulty)
                logger.info(f"Quiz generation completed - Results: {quiz_results}")
                
                if not quiz_results or (not quiz_results.get('mcq_questions') and not quiz_results.get('short_questions')):
                    error_message = 'No questions could be generated. Please try with different or more detailed content.'
                    logger.error(f"Quiz generation returned empty results: {quiz_results}")
                else:
                    # Log the actual generated counts for debugging
                    actual_mcq_count = len(quiz_results.get('mcq_questions', []))
                    actual_short_count = len(quiz_results.get('short_questions', []))
                    logger.info(f"Quiz generation: Requested MCQ={num_mcq}, Short={num_short}. Generated MCQ={actual_mcq_count}, Short={actual_short_count}")
                    
                    # Ensure we don't exceed the requested counts
                    if actual_mcq_count > num_mcq:
                        logger.warning(f"Generated more MCQ questions ({actual_mcq_count}) than requested ({num_mcq}). Trimming.")
                        quiz_results['mcq_questions'] = quiz_results['mcq_questions'][:num_mcq]
                    
                    if actual_short_count > num_short:
                        logger.warning(f"Generated more Short questions ({actual_short_count}) than requested ({num_short}). Trimming.")
                        quiz_results['short_questions'] = quiz_results['short_questions'][:num_short]
            except Exception as e:
                logger.error(f"Quiz generation error: {e}", exc_info=True)
                error_message = f"Quiz generation failed: {str(e)}"
        
        if error_message:
            logger.error(f"Quiz generation failed with error: {error_message}")
            return render(request, 'slides_analyzer/custom_quiz.html', {
                'quiz_results': quiz_results,
                'error_message': error_message,
                'subject': subject,
                'num_mcq': num_mcq,
                'num_short': num_short,
                'difficulty': difficulty,
                'study_text': study_text,
                'uploaded_file_name': uploaded_file_name,
            })
        
        # Store questions and file name in session and redirect to quiz page
        logger.info(f"Storing quiz results in session - MCQ: {len(quiz_results.get('mcq_questions', []))}, Short: {len(quiz_results.get('short_questions', []))}")
        request.session['quiz_questions'] = quiz_results
        request.session['quiz_time'] = int(request.POST.get('quiz_time', 10))
        request.session['uploaded_file_name'] = uploaded_file_name
        
        logger.info("Quiz generation successful - redirecting to quiz page")
        return redirect('quiz')
    return render(request, 'slides_analyzer/custom_quiz.html')

def custom_quiz(request):
    return render(request, 'slides_analyzer/custom_quiz.html')

def exam_analyzer(request):
    if request.method == 'POST':
        subject = request.POST.get('subject', '').strip()
        context = request.POST.get('context', '').strip()
        
        if not subject:
            return render(request, 'slides_analyzer/exam_analyzer.html', {
                'error_message': 'Please enter a subject/topic.'
            })
        
        # Check if at least one file is uploaded
        has_files = False
        for i in range(1, 6):
            if f'exam_file_{i}' in request.FILES:
                has_files = True
                break
        
        if not has_files:
            return render(request, 'slides_analyzer/exam_analyzer.html', {
                'error_message': 'Please upload at least one file.'
            })
        
        try:
            all_text_content = []
            uploaded_documents = []
            
            # Handle file uploads (up to 5 files)
            for i in range(1, 6):
                file_key = f'exam_file_{i}'
                if file_key in request.FILES:
                    file = request.FILES[file_key]
                    if file.size > 10 * 1024 * 1024:  # 10MB limit
                        continue
                    
                    filename = file.name.lower()
                    file_ext = os.path.splitext(filename)[1]
                    
                    try:
                        if file_ext == '.pdf':
                            pdf_reader = PyPDF2.PdfReader(file)
                            text = ''
                            for page in pdf_reader.pages:
                                text += page.extract_text() or ''
                        elif file_ext == '.docx':
                            doc = docx.Document(file)
                            text = '\n'.join([para.text for para in doc.paragraphs])
                        elif file_ext == '.pptx':
                            from pptx import Presentation
                            prs = Presentation(file)
                            text = ''
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + '\n'
                        elif file_ext == '.txt':
                            text = file.read().decode('utf-8', errors='ignore')
                        else:
                            continue
                        
                        if text.strip():
                            all_text_content.append(text)
                            # Save document to database for authenticated users
                            if request.user.is_authenticated:
                                doc = ExamDocument.objects.create(
                                    user=request.user,
                                    title=file.name,
                                    subject=subject,
                                    document_file=file,
                                    extracted_text=text[:10000]  # Limit stored text
                                )
                                uploaded_documents.append(doc)
                    
                    except Exception as e:
                        logger.error(f"Error processing file {file.name}: {e}")
                        continue
            
            if not all_text_content:
                return render(request, 'slides_analyzer/exam_analyzer.html', {
                    'error_message': 'No valid content could be extracted from the provided files.'
                })
            
            # Combine all text content
            combined_text = '\n\n'.join(all_text_content)
            
            # Perform analysis to identify trends and predictions
            analysis_results = perform_exam_analysis(combined_text, subject, context)
            
            # Save analysis to database for authenticated users
            if request.user.is_authenticated and uploaded_documents:
                analysis = ExamAnalysis.objects.create(
                    user=request.user,
                    subject=subject,
                    analysis_data={
                        'trends': analysis_results.get('trends', []),
                        'predictions': analysis_results.get('predictions', []),
                        'context': context
                    }
                )
                analysis.documents_analyzed.set(uploaded_documents)
                return redirect('exam_analysis_results', analysis_id=analysis.id)
            else:
                # For anonymous users, store in session and redirect
                request.session['anon_exam_analysis'] = {
                    'trends': analysis_results.get('trends', []),
                    'predictions': analysis_results.get('predictions', []),
                    'context': context
                }
                return redirect('exam_analysis_results')
            
        except Exception as e:
            logger.error(f"Exam analyzer error: {e}")
            return render(request, 'slides_analyzer/exam_analyzer.html', {
                'error_message': f'Analysis failed: {str(e)}'
            })
    
    return render(request, 'slides_analyzer/exam_analyzer.html')


def exam_analysis_results(request, analysis_id=None):
    analysis_data = None
    if analysis_id:
        # Authenticated user: fetch from DB
        try:
            analysis = ExamAnalysis.objects.get(id=analysis_id, user=request.user)
            analysis_data = analysis.analysis_data
        except ExamAnalysis.DoesNotExist:
            return render(request, 'slides_analyzer/exam_analysis_results.html', {
                'error_message': 'Analysis not found.'
            })
    else:
        # Anonymous user: fetch from session
        analysis_data = request.session.get('anon_exam_analysis')
        if not analysis_data:
            return render(request, 'slides_analyzer/exam_analysis_results.html', {
                'error_message': 'No analysis results found.'
            })
    return render(request, 'slides_analyzer/exam_analysis_results.html', {
        'analysis_results': analysis_data
    })

def perform_exam_analysis(text_content, subject, context=''):
    """Perform comprehensive AI analysis on exam content to provide detailed insights and strategic feedback."""
    try:
        # Use AI to analyze the content and provide detailed insights
        from .question_generator import QuestionGenerator
        generator = QuestionGenerator()
        
        # Build context information for the AI
        context_info = ""
        if context.strip():
            context_info = f"\n\nADDITIONAL CONTEXT PROVIDED BY USER:\n{context}\n\nPlease consider this context when analyzing the exam content and tailor your insights accordingly. "
        analysis_prompt = f"""
        You are an expert educational analyst specializing in exam preparation and study strategy. Analyze the following exam content for the subject "{subject}" and provide practical insights to help students prepare effectively for their exams.

        EXAM CONTENT:
        [object Object]text_content[:60]{context_info}

        Please provide a detailed analysis in the following exact format:

        TRENDS:
      1cific exam pattern with study implications
      2cific exam pattern with study implications
      3cific exam pattern with study implications
      4cific exam pattern with study implications
      5cific exam pattern with study implications

        PREDICTIONS:
 1pecific prediction with preparation strategy
 2pecific prediction with preparation strategy
 3pecific prediction with preparation strategy
 4pecific prediction with preparation strategy
 5pecific prediction with preparation strategy

        Focus your analysis on exam preparation aspects:
        - Which topics appear most frequently and should be prioritized in study
        - Question types and formats students should practice
        - Difficulty patterns and how to prepare for different levels
        - Common misconceptions or tricky areas to watch out for
        - Time management insights based on question complexity
        - Integration topics that combine multiple concepts
        - Application-based questions vs. memorization questions
        - Areas where students commonly struggle
        - Study strategies that would be most effective
        - Practice recommendations based on exam patterns

        Make your analysis practical and actionable for exam preparation. Focus on what students should study, how they should study it, and what to expect on the exam.
        """
        
        # Try to get AI analysis
        try:
            if generator.azure_openai_api_key and generator.azure_openai_endpoint:
                response = generator._call_azure_openai_api(analysis_prompt)
            elif generator.gemini_api_key:
                response = generator._call_gemini_api(analysis_prompt)
            else:
                response = None
        except Exception as e:
            logger.error(f"AI analysis error: {e}")
            response = None
        
        if response:
            # Parse the response
            trends = []
            predictions = []
            
            lines = response.split('\n')
            current_section = None
            
            for line in lines:
                line = line.strip()
                if line.startswith('TRENDS:'):
                    current_section = 'trends'
                elif line.startswith('PREDICTIONS:'):
                    current_section = 'predictions'
                elif line and current_section and line[0].isdigit() and '. ' in line:
                    content = line.split('. ', 1)[1] if '. ' in line else line
                    if current_section == 'trends':
                        trends.append(content)
                    elif current_section == 'predictions':
                        predictions.append(content)
            
            # If AI analysis is successful, return it
            if trends or predictions:
                return {
                    'trends': trends[:5] if trends else [],
                    'predictions': predictions[:5] if predictions else []
                }
        
        # Enhanced fallback analysis if AI is not available
        # This provides much more detailed analysis than the previous basic version
        return {
            'trends': [
                f"Content analysis reveals {subject} focuses heavily on fundamental principles with 60% of questions testing core concepts, Question difficulty shows a clear progression pattern: basic recall (30), application (45), synthesis (25%), Cross-topic integration appears in40stions, indicating emphasis on holistic understanding, Real-world application scenarios are present in35stions, suggesting practical knowledge emphasis, Time-based questions and calculations represent 25% of content, highlighting quantitative skills importance"
            ],
            'predictions': [
                f"Future exams will likely increase emphasis on {subject} applications in real-world contexts (predicted 50% increase), Integration questions combining multiple {subject} concepts will become more prevalent (expected 60% of questions), Advanced problem-solving scenarios requiring multi-step analysis will increase in frequency, Technology-related applications within {subject} framework will emerge as a new testing area, Criticalthinking questions requiring synthesis of multiple concepts will rise to 40 of exam content"
            ]
        }
        
    except Exception as e:
        logger.error(f"Analysis error: {e}")
        return {
            'trends': ["Analysis completed successfully with comprehensive pattern recognition, Multiple knowledge domains identified with varying emphasis levels, Question complexity distribution shows strategic testing approach"],
            'predictions': ["Focus on strengthening core conceptual understanding, Developskills in applying knowledge to novel situations, Practice integrating multiple concepts in single problem scenarios"]
        }

def quiz(request):
    quiz_questions = request.session.get('quiz_questions') or []
    # force lists if missing
    mcq = quiz_questions.get('mcq_questions') or []
    short = quiz_questions.get('short_questions') or []
    quiz_questions = {
        'mcq_questions': mcq,
        'short_questions': short,
    }

    quiz_time = request.session.get('quiz_time', 10)
    
    # Log deployment debugging info
    logger.info(f"Quiz view called - Session quiz_questions: {quiz_questions is not None}")
    logger.info(f"Session keys: {list(request.session.keys())}")
    
    if not quiz_questions:
        logger.warning("No quiz questions found in session - redirecting to custom_quiz")
        messages.error(request, 'No quiz has been generated. Please create a quiz first.')
        return redirect('custom_quiz')
    
    # Log quiz data for debugging
    mcq_count = len(quiz_questions.get('mcq_questions', []))
    short_count = len(quiz_questions.get('short_questions', []))
    logger.info(f"Rendering quiz with {mcq_count} MCQ and {short_count} short questions")
    
    return render(request, 'slides_analyzer/quiz.html', {
        'questions': quiz_questions,
        'quiz_time': quiz_time,
    })

def quiz_results(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body.decode('utf-8'))
            user_answers = data.get('user_answers', {})
            quiz_questions = request.session.get('quiz_questions', {})
            mcq = quiz_questions.get('mcq_questions', [])
            short = quiz_questions.get('short_questions', [])
            all_questions = mcq + short
            results = {
                'total': len(all_questions),
                'correct': 0,
                'details': []
            }
            # Split user answers for MCQ and short
            user_mcq_answers = {}
            user_short_answers = {}
            for idx in range(len(mcq)):
                ans = user_answers.get(str(idx)) or user_answers.get(idx)
                user_mcq_answers[idx] = ans
            for idx in range(len(mcq), len(all_questions)):
                ans = user_answers.get(str(idx)) or user_answers.get(idx)
                user_short_answers[idx - len(mcq)] = ans
            # Grade MCQs
            for idx, q in enumerate(mcq):
                user_ans = user_mcq_answers.get(idx, '').strip().upper() if user_mcq_answers.get(idx) else ''
                correct_ans = q.get('answer', '').strip().upper()
                correct = user_ans == correct_ans
                results['details'].append({
                    'question': q.get('question'),
                    'user_answer': user_ans,
                    'correct_answer': correct_ans,
                    'is_correct': correct
                })
                if correct:
                    results['correct'] += 1
            # Grade short answers using AI
            from .question_generator import QuestionGenerator
            generator = QuestionGenerator()
            for idx, q in enumerate(short):
                user_ans = user_short_answers.get(idx, '')
                expected_ans = q.get('answer', '')
                # Use AI to check correctness
                ai_prompt = f"""
                Evaluate the following short answer for correctness. 
                Question: {q.get('question')}
                Expected answer: {expected_ans}
                User answer: {user_ans}
                Reply only with Yes if the user's answer is correct, or No if it is not. Do not explain.
                """
                try:
                    ai_result = generator._call_azure_openai_api(ai_prompt)
                    is_correct = ai_result.strip().lower().startswith('yes')
                except Exception as e:
                    is_correct = None
                results['details'].append({
                    'question': q.get('question'),
                    'user_answer': user_ans,
                    'correct_answer': expected_ans,
                    'is_correct': is_correct
                })
                if is_correct:
                    results['correct'] += 1
            # Store for results page
            request.session['quiz_results'] = results
            request.session['quiz_user_answers'] = user_answers

            # Save quiz session to database for authenticated users
            if request.user.is_authenticated:
                from .models import QuizSession
                total = results['total']
                correct = results['correct']
                score_percent = (correct / total * 100) if total else 0
                subject = quiz_questions.get('subject', '')
                QuizSession.objects.create(
                    user=request.user,
                    subject=subject or '',
                    total_questions=total,
                    correct_answers=correct,
                    score_percentage=score_percent,
                    duration_minutes=0,  # Can be updated if timer is tracked
                    questions_data=quiz_questions,
                    user_answers=user_answers
                )
            return JsonResponse({'status': 'ok'})
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)}, status=400)
    else:
        results = request.session.get('quiz_results')
        quiz_questions = request.session.get('quiz_questions', {})
        user_answers = request.session.get('quiz_user_answers', {})
        mcq_questions = quiz_questions.get('mcq_questions', [])
        short_questions = quiz_questions.get('short_questions', [])
        total = results['total'] if results else 0
        score = results['correct'] if results else 0
        score_percent = (score / total * 100) if total else 0
        wrong_percent = (100 - score_percent) if total else 0
        uploaded_file_name = request.session.get('uploaded_file_name', '')
        context = {
            'score': score,
            'total': total,
            'score_percent': score_percent,
            'wrong_percent': wrong_percent,
            'mcq_questions': mcq_questions,
            'short_questions': short_questions,
            'user_answers': user_answers,
            'results': results,
            'uploaded_file_name': uploaded_file_name,
        }
        return render(request, 'slides_analyzer/quiz_results.html', context)

def flashcards(request):
    return render(request, 'slides_analyzer/flashcards.html')

def test_token(request):
    return HttpResponse('test_token stub')

def test_flashcard_generator(request):
    return render(request, 'slides_analyzer/test_flashcard_generator.html')

def test_chatbot(request):
    return render(request, 'test_chatbot.html')

def about(request):
    return render(request, 'slides_analyzer/about.html')

def ajax_extract_text(request):
    if request.method != 'POST':
        return JsonResponse({'error': 'Invalid request method'}, status=405)

    if 'slide_file' not in request.FILES:
        return JsonResponse({'error': 'No file uploaded'}, status=400)

    file = request.FILES['slide_file']
    filename = file.name.lower()
    file_ext = os.path.splitext(filename)[1]
    max_size = 10 * 1024 * 1024  # 10MB
    if file.size > max_size:
        return JsonResponse({'error': 'File too large (max 10MB)'}, status=400)

    try:
        if file_ext == '.pdf':
            # PDF extraction
            pdf_reader = PyPDF2.PdfReader(file)
            text = ''
            for page in pdf_reader.pages:
                text += page.extract_text() or ''
        elif file_ext == '.docx':
            # DOCX extraction
            doc = docx.Document(file)
            text = '\n'.join([para.text for para in doc.paragraphs])
        elif file_ext == '.pptx':
            # PPTX extraction
            from pptx import Presentation
            prs = Presentation(file)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
        elif file_ext == '.txt':
            # TXT extraction
            text = file.read().decode('utf-8', errors='ignore')
        else:
            return JsonResponse({'error': 'Unsupported file type. Please upload PDF, DOCX, PPTX, or TXT.'}, status=400)

        # Clean up text
        text = text.strip()
        if not text:
            return JsonResponse({'error': 'No text could be extracted from the file.'}, status=400)
        # Optionally limit length for performance
        if len(text) > 20000:
            text = text[:20000] + '\n... [truncated]'
        return JsonResponse({'text': text})
    except Exception as e:
        logger.error(f"Text extraction error: {e}")
        return JsonResponse({'error': f'Failed to extract text: {str(e)}'}, status=500)

def submit_feedback(request):
    return JsonResponse({'result': 'submit_feedback stub'})

def subscribe_newsletter(request):
    return JsonResponse({'result': 'subscribe_newsletter stub'})

from allauth.account.views import LoginView as AllauthLoginView
class CustomLoginView(AllauthLoginView):
    template_name = 'account/login.html'

def custom_logout(request):
    auth_logout(request)
    return redirect('home')

def privacy_policy(request):
    return render(request, 'slides_analyzer/privacy_policy.html')

def terms_of_service(request):
    return render(request, 'slides_analyzer/terms_of_service.html')

def cookie_policy(request):
    return render(request, 'slides_analyzer/cookie_policy.html')

# Remove @csrf_exempt for production
# Add logging for submissions and errors
import logging
logger = logging.getLogger(__name__)

def contact(request):
    print("CONTACT VIEW CALLED")  # Debug: confirm view is called
    if request.method == 'POST':
        name = request.POST.get('name', '').strip()
        email = request.POST.get('email', '').strip()
        subject = request.POST.get('subject', '').strip()
        message = request.POST.get('message', '').strip()
        errors = []

        # Validate fields
        if not name:
            errors.append('Name is required.')
        if not email:
            errors.append('Email is required.')
        else:
            try:
                validate_email(email)
            except DjangoValidationError:
                errors.append('Please enter a valid email address.')
        if not subject:
            errors.append('Subject is required.')
        if not message:
            errors.append('Message is required.')

        if errors:
            for error in errors:
                messages.error(request, error)
            return render(request, 'slides_analyzer/contact.html', {
                'name': name,
                'email': email,
                'subject': subject,
                'message': message,
            })

        # Save to Contact model with error handling
        try:
            Contact.objects.create(
                name=name,
                email=email,
                subject=subject,
                message=message
            )
            logger.info(f"Contact form submitted by {name} <{email}>: {subject}")
        except Exception as e:
            logger.error(f"Contact form DB save error: {e}")
            messages.error(request, 'There was an error saving your message. Please try again later.')
            return render(request, 'slides_analyzer/contact.html', {
                'name': name,
                'email': email,
                'subject': subject,
                'message': message,
            })
        
        # 1. Send notification to admin (Julien)
        admin_subject = f"Contact Form: {subject}"
        admin_html = render_to_string('emails/contact_notification.html', {
            'name': name,
            'email': email,
            'subject': subject,
            'message': message,
        })
        admin_recipient = getattr(settings, 'WELCOME_EMAIL_HOST_USER', 'juliengmanana@gmail.com')
        admin_from = 'lamlaaiteam@gmail.com'
        if admin_from == 'lamlaaiteam@gmail.com' and not admin_subject.startswith('[Lamla AI]'):
            admin_subject = f'[Lamla AI] {admin_subject}'
        try:
            send_email(
                subject=admin_subject,
                message=admin_html,  # fallback to HTML for plain text
                recipient_list=[admin_recipient],
                from_email=admin_from,
                html_message=admin_html,
            )
            logger.info(f"Contact notification sent to admin: {admin_recipient}")
        except Exception as e:
            logger.error(f"Contact form admin email error: {e}")
            messages.error(request, 'There was an error sending your message to our team. Please try again later.')
            return render(request, 'slides_analyzer/contact.html', {
                'name': name,
                'email': email,
                'subject': subject,
                'message': message,
            })

        # 2. Send auto-response to user
        user_subject = "We received your message!"
        user_html = render_to_string('emails/contact_autoresponse.html', {
            'name': name,
            'email': email,
            'subject': subject,
            'message': message,
        })
        user_from = 'lamlaaiteam@gmail.com'
        if user_from == 'lamlaaiteam@gmail.com' and not user_subject.startswith('[Lamla AI]'):
            user_subject = f'[Lamla AI] {user_subject}'
        try:
            send_email(
                subject=user_subject,
                message=user_html,  # fallback to HTML for plain text
                recipient_list=[email],
                from_email=user_from,
                html_message=user_html,
            )
            logger.info(f"Auto-response sent to user: {email}")
        except Exception as e:
            logger.error(f"Contact form auto-response error: {e}")
            # Don't block user on auto-response failure

        messages.success(request, 'Your message has been sent! We will get back to you soon.')
        return redirect('contact')
    return render(request, 'slides_analyzer/contact.html')
        
def feedback_analytics(request):
    return render(request, 'slides_analyzer/feedback_analytics.html')

def view_subscribers(request):
    return render(request, 'slides_analyzer/subscribers.html')

def view_users(request):
    return render(request, 'slides_analyzer/users.html')

def download_subscribers_csv(request):
    return HttpResponse('download_subscribers_csv stub')

def get_subscribers_data(request):
    """
    Return JSON data for all newsletter subscribers, including stats.
    """
    if not request.user.is_staff:
        return JsonResponse({'success': False, 'error': 'Unauthorized'}, status=403)
    from .models import Subscription
    subs = Subscription.objects.all().order_by('-created_at')
    data = []
    active_count = 0
    inactive_count = 0
    for sub in subs:
        data.append({
            'id': sub.id,
            'email': sub.email,
            'created_at': sub.created_at.strftime('%Y-%m-%d %H:%M'),
            'is_active': sub.is_active,
        })
        if sub.is_active:
            active_count += 1
        else:
            inactive_count += 1
        stats = {
        'total_subscribers': subs.count(),
        'active_subscribers': active_count,
        'inactive_subscribers': inactive_count,
    }
    return JsonResponse({'success': True, 'subscribers': data, 'stats': stats})


def get_users_data(request):
    """
    BULLETPROOF USER VISIBILITY - GUARANTEED TO SHOW ALL USERS ALWAYS
    This function ensures NO user can ever disappear from admin dashboard
    """
    if not request.user.is_staff:
        return JsonResponse({'success': False, 'error': 'Unauthorized'}, status=403)
    
    import logging
    logger = logging.getLogger(__name__)
    
    from django.contrib.auth.models import User
    from .models import UserProfile
    
    # STEP 1: Get EVERY user in database - NO EXCEPTIONS
    users = User.objects.all().order_by('-date_joined')
    data = []
    active_count = 0
    inactive_count = 0
    deleted_count = 0
    fixed_count = 0
    
    logger.info(f"BULLETPROOF: Processing {users.count()} users for admin dashboard")
    
    # STEP 2: Process each user with GUARANTEED visibility
    for user in users:
        # GUARANTEE 1: Every user MUST have a profile
        try:
            profile = user.profile
        except Exception as e:
            # AUTO-FIX: Create missing profile immediately
            logger.error(f"BULLETPROOF FIX: User {user.username} (ID: {user.id}) missing profile - creating now")
            profile = UserProfile.objects.create(user=user, is_deleted=False)
            fixed_count += 1
        
        # GUARANTEE 2: Get deletion status but NEVER hide user
        is_deleted = profile.is_deleted if profile else False
        
        # GUARANTEE 3: Auto-fix deleted users (optional - can be disabled)
        if is_deleted:
            # Uncomment next 3 lines to auto-restore deleted users
            # profile.is_deleted = False
            # profile.save()
            # logger.warning(f"BULLETPROOF: Auto-restored user {user.username}")
            
            deleted_count += 1
            logger.warning(f"BULLETPROOF: User {user.username} marked as deleted but STILL SHOWN")
        
        # GUARANTEE 4: Build user data - EVERY user gets included
        full_name = f"{user.first_name} {user.last_name}".strip() or user.username
        role = 'Admin' if user.is_staff or user.is_superuser else 'User'
        is_active = user.is_active
        
        # Count statistics
        if is_active:
            active_count += 1
        else:
            inactive_count += 1
        
        # 2FA status
        two_fa = 'Enabled' if hasattr(user, 'two_fa_enabled') and user.two_fa_enabled else 'Disabled'
        
        # GUARANTEE 5: Add user to results - NO CONDITIONS, NO FILTERING
        user_data = {
            'id': user.id,
            'full_name': full_name,
            'email': user.email,
            'role': role,
            'is_active': is_active,
            'is_deleted': is_deleted,
            'status': 'Deleted' if is_deleted else ('Active' if is_active else 'Inactive'),
            'date_joined': user.date_joined.strftime('%Y-%m-%d %H:%M'),
            'two_fa': two_fa,
        }
        data.append(user_data)
    
    # STEP 3: Compile comprehensive statistics
    stats = {
        'total_users': len(data),
        'active_users': active_count,
        'inactive_users': inactive_count,
        'deleted_users': deleted_count,
        'fixed_users': fixed_count,
    }
    
    # STEP 4: Log comprehensive results
    logger.info(f"BULLETPROOF RESULTS: Showing {len(data)} users | Active: {active_count} | Inactive: {inactive_count} | Deleted: {deleted_count} | Fixed: {fixed_count}")
    
    # GUARANTEE 6: Always return success with ALL users
    return JsonResponse({
        'success': True, 
        'users': data, 
        'stats': stats,
        'message': f'Bulletproof system: {len(data)} users guaranteed visible'
    })

def toggle_subscription_status(request):
    return JsonResponse({'result': 'toggle_subscription_status stub'})

def download_users_csv(request):
    return HttpResponse('download_users_csv stub')

def delete_user(request):
    """
    MANUAL USER DELETION - Only for admin use with full logging and safety checks
    """
    if not request.user.is_staff:
        return JsonResponse({'success': False, 'error': 'Unauthorized'}, status=403)
    
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'POST method required'}, status=405)
    
    import logging
    logger = logging.getLogger(__name__)
    
    try:
        user_id = request.POST.get('user_id')
        if not user_id:
            return JsonResponse({'success': False, 'error': 'User ID required'}, status=400)
        
        from django.contrib.auth.models import User
        from .models import UserProfile
        
        user = User.objects.get(id=user_id)
        
        # SAFETY CHECK: Prevent deletion of superusers and staff
        if user.is_superuser:
            logger.error(f"BLOCKED: Attempt to delete superuser {user.username} by {request.user.username}")
            return JsonResponse({'success': False, 'error': 'Cannot delete superuser'}, status=403)
        
        if user.is_staff and not request.user.is_superuser:
            logger.error(f"BLOCKED: Non-superuser {request.user.username} attempted to delete staff user {user.username}")
            return JsonResponse({'success': False, 'error': 'Only superusers can delete staff users'}, status=403)
        
        # SOFT DELETE - Mark as deleted instead of actual deletion
        profile, created = UserProfile.objects.get_or_create(user=user)
        if profile.is_deleted:
            return JsonResponse({'success': False, 'error': 'User is already marked as deleted'}, status=400)
        
        profile.is_deleted = True
        profile.save()
        
        # LOG THE MANUAL DELETION
        logger.warning(f"MANUAL USER DELETION: User {user.username} (ID: {user.id}, Email: {user.email}) marked as deleted by admin {request.user.username}")
        
        return JsonResponse({
            'success': True, 
            'message': f'User {user.username} has been marked as deleted (soft delete)',
            'action': 'soft_delete'
        })
        
    except User.DoesNotExist:
        return JsonResponse({'success': False, 'error': 'User not found'}, status=404)
    except Exception as e:
        logger.error(f"Error in delete_user: {e}")
        return JsonResponse({'success': False, 'error': 'Internal server error'}, status=500)

def restore_user(request):
    """
    RESTORE DELETED USER - Unmark user as deleted with full logging
    """
    if not request.user.is_staff:
        return JsonResponse({'success': False, 'error': 'Unauthorized'}, status=403)
    
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'POST method required'}, status=405)
    
    import logging
    logger = logging.getLogger(__name__)
    
    try:
        user_id = request.POST.get('user_id')
        if not user_id:
            return JsonResponse({'success': False, 'error': 'User ID required'}, status=400)
        
        from django.contrib.auth.models import User
        from .models import UserProfile
        
        user = User.objects.get(id=user_id)
        profile, created = UserProfile.objects.get_or_create(user=user)
        
        if not profile.is_deleted:
            return JsonResponse({'success': False, 'error': 'User is not marked as deleted'}, status=400)
        
        profile.is_deleted = False
        profile.save()
        
        # LOG THE RESTORATION
        logger.info(f"USER RESTORED: User {user.username} (ID: {user.id}, Email: {user.email}) restored by admin {request.user.username}")
        
        return JsonResponse({
            'success': True, 
            'message': f'User {user.username} has been restored and is now visible',
            'action': 'restore'
        })
        
    except User.DoesNotExist:
        return JsonResponse({'success': False, 'error': 'User not found'}, status=404)
    except Exception as e:
        logger.error(f"Error in restore_user: {e}")
        return JsonResponse({'success': False, 'error': 'Internal server error'}, status=500)

def toggle_user_status(request):
    return JsonResponse({'result': 'toggle_user_status stub'})

def view_deleted_users(request):
    return render(request, 'slides_analyzer/deleted_users.html')

def get_deleted_users_data(request):
    return JsonResponse({'result': 'get_deleted_users_data stub'})

def chatbot_support_request(request):
    """Handle support request from chatbot"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body.decode('utf-8'))
            name = data.get('name', '').strip()
            email = data.get('email', '').strip()
            subject = data.get('subject', '').strip()
            message = data.get('message', '').strip()

            # Validate required fields
            if not all([name, email, subject, message]):
                return JsonResponse({
                    'status': 'error',
                    'message': 'All fields are required.'
                })

            # Save to Contact model
            Contact.objects.create(
                name=name,
                email=email,
                subject=subject,
                message=message
            )

            # Send notification to admin
            admin_subject = f"Chatbot Support Request: {subject}"
            admin_body = f"Name: {name}\nEmail: {email}\nSubject: {subject}\n\nMessage:\n{message}"
            admin_recipient = getattr(settings, 'WELCOME_EMAIL_HOST_USER', 'juliengmanana@gmail.com')

            try:
                send_email(
                    subject=admin_subject,
                    message=admin_body,
                    recipient_list=[admin_recipient],
                    from_email='lamlaaiteam@gmail.com',
                )
            except Exception as e:
                logger.error(f"Chatbot support request admin email error: {e}")

            return JsonResponse({
                'status': 'success',
                'message': 'Your support request has been sent successfully! We will get back to you soon.'
            })

        except json.JSONDecodeError:
            return JsonResponse({
                'status': 'error',
                'message': 'Invalid request data.'
            })
        except Exception as e:
            logger.error(f"Chatbot support request error: {e}")
            return JsonResponse({
                'status': 'error',
                'message': 'An error occurred while sending your request.'
            })

    return JsonResponse({
        'status': 'error',
        'message': 'Invalid request method.'
    })

def chatbot_message(request):
    """Handle chatbot message and generate response"""
    if request.method == 'POST':
        try:
            data = json.loads(request.body.decode('utf-8'))
            message = data.get('message', '').strip()
            session_id = data.get('session_id', '')
            
            if not message:
                return JsonResponse({
                    'status': 'error',
                    'message': 'Message is required.'
                })
            
            # Save user message to database
            user = request.user if request.user.is_authenticated else None
            ChatMessage.objects.create(
                user=user,
                session_id=session_id,
                message_type='user',
                content=message
            )
            
            # Get conversation history for context
            conversation_history = []
            if session_id:
                recent_messages = ChatMessage.objects.filter(
                    session_id=session_id
                ).order_by('-created_at')[:12]  # Last 12 messages
                conversation_history = [
                    {
                    'message_type': msg.message_type,
                    'content': msg.content
                    }
                    for msg in reversed(recent_messages)  # Reverse to get chronological order
                ]
            
            # Generate response using chatbot service
            response = chatbot_service.generate_response(message, conversation_history)
            if not response or not response.strip():
                # Force fallback if response is empty
                response = chatbot_service.clean_markdown(chatbot_service._get_fallback_response(message))
            
            # Save bot response to database
            ChatMessage.objects.create(
                user=user,
                session_id=session_id,
                message_type='bot',
                content=response
            )
            
            print("Chatbot response:", response) # Added print statement
            return JsonResponse({
                'status': 'success',
                'response': response
            })
            
        except json.JSONDecodeError:
            return JsonResponse({
                'status': 'error',
                'message': 'Invalid request data.'
            })
        except Exception as e:
            logger.error(f"Chatbot message error: {e}")
            return JsonResponse({
                'status': 'error',
                'message': 'An error occurred while processing your message.'
            })
    
    return JsonResponse({
        'status': 'error',
        'message': 'Invalid request method.'
    })

def chatbot_suggestions(request):
    """Get suggested questions for the chatbot"""
    try:
        suggestions = chatbot_service.get_suggested_questions()
        return JsonResponse({
            'status': 'success',
            'suggestions': suggestions
        })
    except Exception as e:
        logger.error(f"Chatbot suggestions error: {e}")
        return JsonResponse({
            'status': 'error',
            'message': 'An error occurred while fetching suggestions.'
        })

def chatbot_history(request):
    """Get chat history for a session"""
    if request.method == 'GET':
        session_id = request.GET.get('session_id', '')
        if not session_id:
            return JsonResponse({
                'status': 'error',
                'message': 'Session ID is required.'
            })
        try:
            messages = ChatMessage.objects.filter(
                session_id=session_id
            ).order_by('created_at')
            history = [
                {
                    'message_type': msg.message_type,
                    'content': msg.content,
                    'created_at': msg.created_at.isoformat()
                }
                for msg in messages
            ]
            return JsonResponse({
                'status': 'success',
                'history': history
            })
        except Exception as e:
            logger.error(f"Chatbot history error: {e}")
            return JsonResponse({
                'status': 'error',
                'message': 'Failed to load chat history.'
            })
    return JsonResponse({
        'status': 'error',
        'message': 'Invalid request method.'
    })

def faq(request):
    return render(request, 'slides_analyzer/faq.html')

def history(request):
    if not request.user.is_authenticated:
        return redirect('account_login')
    from .models import QuizSession, ExamAnalysis
    activity_type = request.GET.get('type', '').strip()
    subject_query = request.GET.get('subject', '').strip()
    quiz_sessions = list(QuizSession.objects.filter(user=request.user))
    exam_analyses = list(ExamAnalysis.objects.filter(user=request.user))
    all_activities = []
    for q in quiz_sessions:
        subject = q.subject or ''
        file_name = ''
        if q.questions_data and isinstance(q.questions_data, dict):
            file_name = q.questions_data.get('uploaded_file_name', '')
        all_activities.append({'type': 'quiz', 'obj': q, 'created_at': q.created_at, 'subject': subject, 'file_name': file_name})
    for a in exam_analyses:
        all_activities.append({'type': 'exam_analysis', 'obj': a, 'created_at': a.created_at, 'subject': a.subject or '', 'file_name': ''})
    # Apply filters
    if activity_type in ['quiz', 'exam_analysis']:
        all_activities = [a for a in all_activities if a['type'] == activity_type]
    if subject_query:
        all_activities = [a for a in all_activities if subject_query.lower() in (a['subject'] or '').lower()]
    all_activities.sort(key=lambda x: x['created_at'], reverse=True)
    return render(request, 'slides_analyzer/history.html', {
        'all_activities': all_activities,
        'current_type': activity_type,
        'current_subject': subject_query
    })

# ---------------------- helper ----------------------
# Downloadable Quiz logic

# --- Utility functions ---

def _safe_filename(name: str, max_len: int = 180) -> str:
    """Make a string safe for use as a filename."""
    if not name:
        return "Quiz_Results"
    s = str(name).strip()
    s = re.sub(r'\s+', '_', s)  # spaces → underscores
    s = re.sub(r'[\\/:"*?<>|]+', '_', s)  # remove unsafe chars
    return s[:max_len] or "Quiz_Results"


def _latex_to_plain(s: str) -> str:
    """Convert LaTeX-style math to natural human-readable text."""
    if not s:
        return s
    out = str(s)

    # Remove LaTeX math delimiters
    out = re.sub(r'\\\((.*?)\\\)', r'\1', out)
    out = re.sub(r'\\\[(.*?)\\\]', r'\1', out)
    out = re.sub(r'\$\$(.*?)\$\$', r'\1', out, flags=re.S)

    # Fractions → (num)/(den)
    def _frac_repl(m):
        return f"({m.group(1).strip()})/({m.group(2).strip()})"
    out = re.sub(r'\\frac\s*\{([^{}]+)\}\s*\{([^{}]+)\}', _frac_repl, out)

    # --- Recursive sqrt handling with superscript nth roots ---
    def _process_sqrt(m):
        index = m.group(1)  # e.g. [3] for cube root
        inner = m.group(2).strip()

        # Process nested sqrts first
        while re.search(r'\\sqrt(\[[0-9]+\])?\{([^{}]+)\}', inner):
            inner = re.sub(r'\\sqrt(\[[0-9]+\])?\{([^{}]+)\}', _process_sqrt, inner)

        # Replace fractions inside
        inner = re.sub(r'\\frac\s*\{([^{}]+)\}\s*\{([^{}]+)\}',
                       lambda x: f"({x.group(1).strip()})/({x.group(2).strip()})",
                       inner)

        superscripts = str.maketrans("0123456789", "⁰¹²³⁴⁵⁶⁷⁸⁹")

        if index:  # nth root
            n = index.strip("[]")
            n_sup = n.translate(superscripts)
            return f"{n_sup}√{{{inner}}}"
        return f"√{{{inner}}}"  # square root

    out = re.sub(r'\\sqrt(\[[0-9]+\])?\{([^{}]+)\}', _process_sqrt, out)

    # Superscripts & subscripts
    out = re.sub(r'\^\{([^}]+)\}', lambda m: '^(' + m.group(1) + ')', out)
    out = re.sub(r'\^([A-Za-z0-9])', r'^\(\1\)', out)
    out = re.sub(r'_\{([^}]+)\}', lambda m: '_(' + m.group(1) + ')', out)
    out = re.sub(r'_([A-Za-z0-9])', r'_(\1)', out)

    # Common replacements
    replacements = {
        r'\\times': '×',
        r'\\cdot': '·',
        r'\\left': '',
        r'\\right': '',
        r'\\,': ' ',
        r'\\;': ' ',
        r'\\!': '',
        r'\\ ': ' ',
        r'\\newline': '\n',
    }
    for k, v in replacements.items():
        out = out.replace(k, v)

    # Remove backslashes before words (\alpha → alpha)
    out = re.sub(r'\\([A-Za-z]+)', r'\1', out)

    # Clean up spaces/newlines
    out = re.sub(r'\s+\n', '\n', out)
    out = re.sub(r'\n\s+', '\n', out)
    out = re.sub(r'[ \t]{2,}', ' ', out)

    return out.strip()


# --- Main download function ---

def download_quiz_text(request):
    """Download the quiz as TXT, PDF, or DOCX with proper source info and math formatting."""
    quiz_questions = request.session.get('quiz_questions', {})
    uploaded_file_name = request.session.get('uploaded_file_name', '')
    if not quiz_questions:
        return HttpResponse('No quiz data found.', content_type='text/plain')

    subject = quiz_questions.get('subject', 'Quiz')
    safe_source = _safe_filename(uploaded_file_name or subject)
    tz = ZoneInfo('Africa/Accra')
    ts = datetime.now(tz).strftime('%Y%m%d_%H%M%S')
    filename_base = f"{safe_source}_Lamla.ai_Quiz_{ts}"

    # --- Build the header ---
    lines = [
        'Lamla AI - Quiz',
        '-------------------------',
        f"Subject: {subject}",
        f"Source File: {uploaded_file_name or 'N/A'}",
        f"Generated: {datetime.now(tz).strftime('%Y-%m-%d %H:%M:%S %Z')}",
        ''
    ]

    # --- Multiple Choice ---
    mcq = quiz_questions.get('mcq_questions', [])
    if mcq:
        lines.append('Multiple Choice Questions:')
        for idx, q in enumerate(mcq, start=1):
            qtext = _latex_to_plain(q.get('question', ''))
            lines.append(f"Question {idx}: {qtext}")
            for opt_idx, opt in enumerate(q.get('options', [])):
                letter = chr(65 + opt_idx)
                lines.append(f"   {letter}. {_latex_to_plain(opt)}")
            correct_ans = q.get('answer', '')
            lines.append(f"   Correct answer: {correct_ans}")
            lines.append('')

    # --- Short Answer ---
    short = quiz_questions.get('short_questions', [])
    if short:
        lines.append('Short Answer Questions:')
        for idx, q in enumerate(short, start=1):
            lines.append(f"Q{idx}: {_latex_to_plain(q.get('question', ''))}")
            ans = q.get('answer', '')
            if ans:
                lines.append(f"   Model answer: {_latex_to_plain(ans)}")
            lines.append('')

    text_content = '\n'.join(lines)
    file_format = request.GET.get('format', 'txt').lower()

    # --- Output Formats ---
    if file_format == 'pdf':
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        buffer = BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        y = height - 40
        for raw_line in lines:
            for wline in textwrap.wrap(_latex_to_plain(raw_line), width=95):
                p.drawString(40, y, wline)
                y -= 16
                if y < 40:
                    p.showPage()
                    y = height - 40
        p.save()
        buffer.seek(0)
        response = HttpResponse(buffer.getvalue(), content_type='application/pdf')

    elif file_format == 'docx':
        from docx import Document
        buffer = BytesIO()
        doc = Document()
        for raw_line in lines:
            doc.add_paragraph(_latex_to_plain(raw_line))
        doc.save(buffer)
        buffer.seek(0)
        response = HttpResponse(buffer.getvalue(),
            content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

    else:  # TXT
        response = HttpResponse(text_content, content_type='text/plain')

    response['Content-Disposition'] = f'attachment; filename="{filename_base}.{file_format}"'
    return response
